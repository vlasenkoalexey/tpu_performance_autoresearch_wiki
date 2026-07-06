#!/usr/bin/env python3
"""Generate an editable PPTX deck for one experiment series (every experiment, full).

    python tools/demo/build_deck_pptx.py --series qwen3-ag-jax --seq 2048 --out /path/deck.pptx

Native, editable slides (python-pptx) that mirror the HTML player's look:
  - the player's exact light-theme palette + agent accent colour;
  - phase emojis (💡 Hypothesis 🔧 Change 📊 Profile 📈 Result ⚖️ Verdict);
  - `code` / **bold** spans rendered as accent highlights (e.g. `bs=16`);
  - GitHub-style red-/green+ diffs with per-line tint;
  - every graph slide is an ANIMATED GIF of the MFU frontier climbing up to that
    experiment (plus a final full-climb slide).
Same extraction as the player, via build_demo. Charts: matplotlib; GIFs: Pillow.
"""
import argparse, os, re, sys, tempfile

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, SCRIPT_DIR)
import build_demo as bd

AGENTS = bd.AGENTS

# ---- palette: EXACTLY the HTML player's light theme (player.html.template) ----
INK  = RGBColor(0x1A, 0x1D, 0x21)   # --text
MUTE = RGBColor(0x5B, 0x65, 0x70)   # --muted
SUBT = RGBColor(0x8A, 0x92, 0x9C)   # --subt
LINE = RGBColor(0xE6, 0xE9, 0xEE)   # --border
GRID = RGBColor(0xEC, 0xEF, 0xF3)   # --grid
CODEBG = RGBColor(0xFB, 0xFC, 0xFD) # --codebg
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ADD_FG = RGBColor(0x11, 0x63, 0x29); ADD_BG = RGBColor(0xE6, 0xFF, 0xED)
DEL_FG = RGBColor(0xA0, 0x13, 0x26); DEL_BG = RGBColor(0xFF, 0xEE, 0xF0)
HUNK = RGBColor(0x6F, 0x42, 0xC1); META = RGBColor(0x8A, 0x92, 0x9C)
VBADGE = {"supported": (ADD_FG, ADD_BG), "refuted": (DEL_FG, DEL_BG),
          "invalid": (DEL_FG, DEL_BG), "inconclusive": (MUTE, GRID),
          "baseline": (MUTE, GRID)}
PHASE_EMOJI = {"Hypothesis": "💡", "Change applied": "🔧", "Profile": "📊",
               "Result": "📈", "Verdict": "⚖️"}
MONO = "Consolas"; SANS = "Calibri"
SW, SH = Inches(13.333), Inches(7.5)


def hexrgb(h):
    h = h.lstrip("#"); return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def no_line(shape):
    shape.line.fill.background()


# ---------------------------------------------------------------- chart / gif
from matplotlib.lines import Line2D

SEQ_LABEL = {2048: "2k", 8192: "8k", 4096: "4k", 1024: "1k"}


def _compress(path, colors=48):
    Image.open(path).convert("RGB").quantize(colors=colors, method=Image.MEDIANCUT).save(path, optimize=True)


def _shade(hex_, f):
    h = hex_.lstrip("#"); r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return "#%02x%02x%02x" % (int(r * f), int(g * f), int(b * f))


def seq_palette(accent, seqs):
    """Per-context colours: first context = accent, others = darker shades of it."""
    shades = [1.0, 0.56, 0.4]
    return {s: (accent if i == 0 else _shade(accent, shades[min(i, 2)])) for i, s in enumerate(seqs)}


def _seq_frontier(exps, upto, s):
    best = None; fx = []; fy = []
    for i in range(min(upto, len(exps) - 1) + 1):
        e = exps[i]
        if e.get("seq") != s:
            continue
        m = e.get("mfu")
        if m is not None and (e.get("verdict") or "").lower() != "invalid":
            best = m if best is None or m > best else best
        if best is not None:
            fx.append(i + 1); fy.append(best)
    return best, fx, fy


def _yrange(exps, targets):
    mfus = [e["mfu"] for e in exps if e.get("mfu") is not None]
    ceil = max(mfus + list(targets.values()) + [10]) * 1.14
    floor = min(mfus + [ceil]) * 0.8
    return floor, ceil


def _decorate(ax, targets, colors, seqs, N, floor, ceil):
    for s in seqs:
        tv = targets.get(s)
        if not tv:
            continue
        ax.axhline(tv, ls=(0, (6, 5)), color=colors[s], alpha=0.55, lw=1.3, zorder=2)
        lbl = (f"MaxText {SEQ_LABEL.get(s, s)} {tv:.1f}%") if len(seqs) > 1 else f"MaxText SOTA {tv:.1f}%"
        ax.text(N + 0.4, tv, lbl, ha="right", va="bottom", color=colors[s], fontsize=10, alpha=0.9)
    ax.set_xlim(0.4, N + 0.6); ax.set_ylim(floor, ceil)
    ax.set_ylabel("MFU (%)", fontsize=11, color="#5b6570")
    ax.set_xlabel("experiment #", fontsize=11, color="#5b6570")
    ax.grid(True, axis="y", color="#eceff3", lw=1)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    for sp in ("left", "bottom"):
        ax.spines[sp].set_color("#d7dbe0")
    ax.tick_params(colors="#8a929c", labelsize=9)
    if len(seqs) > 1:
        handles = [Line2D([0], [0], marker="o", linestyle="", markerfacecolor=colors[s],
                          markeredgecolor=colors[s], label=SEQ_LABEL.get(s, str(s)), markersize=9)
                   for s in seqs]
        leg = ax.legend(handles=handles, loc="lower right", frameon=False, fontsize=11,
                        handletextpad=0.3, labelspacing=0.3, title="context")
        leg.get_title().set_fontsize(9); leg.get_title().set_color("#8a929c")


def _best_label(ax, sk, best, colors, accent, seqs):
    if best is None:
        return
    lab = f"Best {SEQ_LABEL.get(sk, '')} MFU: {best:.1f}%" if len(seqs) > 1 else f"Best MFU: {best:.1f}%"
    ax.text(0.015, 0.965, lab, transform=ax.transAxes, fontsize=24 if len(seqs) > 1 else 27,
            fontweight="bold", color=colors.get(sk, accent), va="top", ha="left")


def draw_chart(exps, k, path, accent, targets, seqs, w=8.9, h=4.7, dpi=95):
    """Full snapshot up to experiment k; points/frontiers coloured per context (2k/8k)."""
    N = len(exps); colors = seq_palette(accent, seqs)
    fig, ax = plt.subplots(figsize=(w, h), dpi=dpi)
    fig.patch.set_facecolor("white"); ax.set_facecolor("white")
    floor, ceil = _yrange(exps, targets)
    for s in seqs:
        _, fx, fy = _seq_frontier(exps, k, s)
        if len(fx) > 1:
            ax.step(fx, fy, where="post", color=colors[s], lw=3, zorder=4, solid_capstyle="round")
    for i in range(k + 1):
        e = exps[i]; m = e.get("mfu")
        if m is None:
            continue
        col = colors.get(e.get("seq"), accent)
        fail = (e.get("verdict") or "").lower() in ("invalid", "refuted", "inconclusive")
        ax.scatter(i + 1, m, s=170 if i == k else 55, facecolors="white" if fail else col,
                   edgecolors=col, linewidths=2.4 if i == k else 1.4, zorder=6)
    sk = exps[k].get("seq"); bestk, _, _ = _seq_frontier(exps, k, sk)
    _best_label(ax, sk, bestk, colors, accent, seqs)
    _decorate(ax, targets, colors, seqs, N, floor, ceil)
    fig.tight_layout(); fig.savefig(path, facecolor="white"); plt.close(fig); _compress(path)


def make_gif(frame_imgs, upto, path, hold=5, dur=470):
    seq = frame_imgs[:upto + 1] + [frame_imgs[upto]] * hold
    seq[0].save(path, save_all=True, append_images=seq[1:], duration=dur, loop=0, optimize=True)


def _ease(t):
    t = max(0.0, min(1.0, t))
    return t * t * (3 - 2 * t)


def draw_frame(exps, k, t, path, accent, targets, seqs, w=8.9, h=4.7, dpi=95):
    """One transition frame: prior state (0..k-1) + experiment k animating in by t.
    The new point drops in and only ITS context's frontier extends one step."""
    N = len(exps); colors = seq_palette(accent, seqs)
    fig, ax = plt.subplots(figsize=(w, h), dpi=dpi)
    fig.patch.set_facecolor("white"); ax.set_facecolor("white")
    floor, ceil = _yrange(exps, targets)
    et = _ease(t); sk = exps[k].get("seq"); colk = colors.get(sk, accent)
    # prior frontiers (each context, up to k-1) + prior points
    for s in seqs:
        _, fx, fy = _seq_frontier(exps, k - 1, s)
        if len(fx) > 1:
            ax.step(fx, fy, where="post", color=colors[s], lw=3, zorder=4, solid_capstyle="round")
    for i in range(k):
        e = exps[i]; m = e.get("mfu")
        if m is None:
            continue
        col = colors.get(e.get("seq"), accent)
        fail = (e.get("verdict") or "").lower() in ("invalid", "refuted", "inconclusive")
        ax.scatter(i + 1, m, s=55, facecolors="white" if fail else col, edgecolors=col,
                   linewidths=1.4, zorder=6)
    # animate experiment k of context sk
    mk = exps[k].get("mfu"); prev, pfx, _ = _seq_frontier(exps, k - 1, sk); newbest = prev
    if mk is not None:
        failk = (exps[k].get("verdict") or "").lower() in ("invalid", "refuted", "inconclusive")
        if (exps[k].get("verdict") or "").lower() != "invalid":
            newbest = mk if prev is None or mk > prev else prev
        if et > 0.01:
            if prev is not None and pfx:
                xprev = pfx[-1]
                ax.plot([xprev, xprev + et * (k + 1 - xprev)], [prev, prev], color=colk, lw=3,
                        solid_capstyle="round", zorder=4)
                if et >= 0.999 and newbest > prev:
                    ax.plot([k + 1, k + 1], [prev, newbest], color=colk, lw=3, zorder=4)
            ax.scatter(k + 1, mk, s=170 * et, facecolors="white" if failk else colk,
                       edgecolors=colk, linewidths=2.4, alpha=min(1.0, et * 1.3), zorder=7)
    disp = newbest if (t >= 0.999 or prev is None) else prev
    _best_label(ax, sk, disp, colors, accent, seqs)
    _decorate(ax, targets, colors, seqs, N, floor, ceil)
    fig.tight_layout(); fig.savefig(path, facecolor="white"); plt.close(fig); _compress(path)


def build_transition_gif(exps, k, path, accent, targets, seqs, tmp):
    """A GIF showing only the single move: prior state → experiment k's state, then hold."""
    ts = [0.0, 0.16, 0.34, 0.52, 0.7, 0.86, 1.0]
    imgs = []
    for j, t in enumerate(ts):
        fp = os.path.join(tmp, f"tr{k:03d}_{j}.png")
        draw_frame(exps, k, t, fp, accent, targets, seqs)
        imgs.append(Image.open(fp).convert("RGB").quantize(colors=48, method=Image.MEDIANCUT))
    durations = [560, 90, 90, 90, 90, 110, 1800]   # hold prior · quick move · hold new
    imgs[0].save(path, save_all=True, append_images=imgs[1:],
                 duration=durations, loop=0, optimize=True)


# ---------------------------------------------------------------- text helpers
def _run(p, text, size, color, bold=False, mono=False):
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold
    f.name = MONO if mono else SANS; f.color.rgb = color
    return r


_RICH = re.compile(r"(`[^`]+`|\*\*[^*]+\*\*|__[^_]+__)")


def add_rich(p, text, size, base, accent):
    """Split text into runs; `code`/**bold**/__x__ → accent-bold highlight (e.g. bs=16)."""
    for tok in _RICH.split(text):
        if not tok:
            continue
        if tok[:1] in "`*_" and len(tok) > 2 and tok[-1:] in "`*_":
            _run(p, tok.strip("`*_"), size, accent, bold=True)
        else:
            _run(p, tok, size, base)


def add_bg(slide, rgb=WHITE):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    return s


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    return tf


def footer(slide, accent, agent_name, n, total):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.02), Inches(12.13), Pt(1))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE; no_line(ln)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(7.16), Inches(0.13), Inches(0.13))
    dot.fill.solid(); dot.fill.fore_color.rgb = accent; no_line(dot)
    tf = textbox(slide, Inches(0.82), Inches(7.08), Inches(9), Inches(0.32))
    _run(tf.paragraphs[0], f"{agent_name} · autonomous TPU optimization", 10.5, MUTE)
    tf2 = textbox(slide, Inches(10.5), Inches(7.08), Inches(2.2), Inches(0.32))
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    _run(p, f"{n} / {total}", 10.5, SUBT)


ICON_PNG = None   # set by main() to the agent's rasterised logo, drawn before the exp title


def agent_icon_png(agent, tmp):
    """Path to a raster logo for the agent (pptx can't embed SVG). Prefer assets/<agent>.png,
    else rasterise assets/<agent>.svg with cairosvg."""
    png = os.path.join(SCRIPT_DIR, "assets", f"{agent}.png")
    if os.path.exists(png):
        return png
    svg = os.path.join(SCRIPT_DIR, "assets", f"{agent}.svg")
    if os.path.exists(svg):
        out = os.path.join(tmp, f"{agent}_icon.png")
        try:                                    # resvg handles masks/blur filters correctly
            import resvg_py
            open(out, "wb").write(bytes(resvg_py.svg_to_bytes(
                svg_string=open(svg, encoding="utf-8").read(), width=1024)))
            return out
        except Exception:
            try:                                # cairosvg fallback (may mangle masked SVGs)
                import cairosvg
                cairosvg.svg2png(url=svg, write_to=out, output_width=512)
                return out
            except Exception:
                return None
    return None


def header(slide, e, phase, accent):
    tf = textbox(slide, Inches(0.6), Inches(0.3), Inches(12.13), Inches(0.95))
    p = tf.paragraphs[0]
    _run(p, f"exp #{e['vnum']} — {e.get('title') or e['slug']}", 23, INK, bold=True)
    if e.get("page_url"):   # link only the arrow so the title stays dark (not hyperlink-blue)
        _run(p, "  ↗", 15, accent, bold=True).hyperlink.address = e["page_url"]
    p2 = tf.add_paragraph(); p2.space_before = Pt(2)
    _run(p2, PHASE_EMOJI.get(phase, "") + "  ", 14, INK)
    _run(p2, phase.upper(), 12.5, accent, bold=True)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.24), Inches(12.13), Pt(1.4))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE; no_line(ln)


def model_first(diff):
    """Split a unified diff into per-file blocks and order model/trainer/sharding files
    BEFORE launch-config (run.sh). Keeps the code change front-and-centre when both exist."""
    if not diff.strip():
        return []
    blocks, cur = [], []
    for ln in diff.split("\n"):
        if (ln.startswith("diff ") or ln.startswith("--- ")) and cur:
            blocks.append(cur); cur = [ln]
        else:
            cur.append(ln)
    if cur:
        blocks.append(cur)
    if len(blocks) <= 1:
        return diff.split("\n")
    def is_config(b):
        head = "\n".join(b[:3]).lower()
        return "run.sh" in head or "launch" in head or "command" in head
    blocks.sort(key=is_config)   # False(model)=0 first, True(config)=1 last
    out = []
    for b in blocks:
        out.extend(b)
    return out


# ---------------------------------------------------------------- slides
def add_hyp(prs, e, accent, name, n, total):
    s = blank(prs); header(s, e, "Hypothesis", accent)
    tf = textbox(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(4.9), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.line_spacing = 1.32
    add_rich(p, e.get("hypothesis_statement") or e.get("hypothesis_title") or "", 21, INK, accent)
    if e.get("expected_gain"):
        p2 = tf.add_paragraph(); p2.space_before = Pt(16)
        _run(p2, "🎯 expected: ", 15, MUTE); add_rich(p2, e["expected_gain"], 15, MUTE, accent)
    footer(s, accent, name, n, total)


def add_change(prs, e, accent, name, n, total):
    s = blank(prs); header(s, e, "Change applied", accent)
    tf = textbox(s, Inches(0.6), Inches(1.62), Inches(12.1), Inches(0.5))
    _run(tf.paragraphs[0], e.get("diff_summary") or "", 15, INK, bold=True)
    kind = {"code": "model code diff", "flag": "launch-config diff",
            "none": "configuration change"}.get(e.get("diff_kind"), "change")
    tf2 = textbox(s, Inches(0.6), Inches(2.02), Inches(12.1), Inches(0.32))
    _run(tf2.paragraphs[0], kind, 11.5, SUBT)
    # diff card + per-line tint — model/trainer files shown BEFORE run.sh (config)
    lines = model_first(e.get("diff") or "")
    MAXL = 24
    more = max(0, len(lines) - MAXL); lines = lines[:MAXL]
    cx, cy, cw = Inches(0.6), Inches(2.45), Inches(12.13)
    lh = Inches(0.185)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, lh * (len(lines) + (1 if more else 0)) + Inches(0.16))
    card.fill.solid(); card.fill.fore_color.rgb = CODEBG; card.line.color.rgb = LINE
    card.adjustments[0] = 0.02
    y = cy + Inches(0.08)
    for ln in lines:
        fg, bg = INK, None
        if ln.startswith("+++") or ln.startswith("---") or ln.startswith("diff "):
            fg = META
        elif ln.startswith("@@"):
            fg = HUNK
        elif ln.startswith("+"):
            fg, bg = ADD_FG, ADD_BG
        elif ln.startswith("-"):
            fg, bg = DEL_FG, DEL_BG
        row = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx + Inches(0.06), y, cw - Inches(0.12), lh)
        no_line(row)
        if bg is not None:
            row.fill.solid(); row.fill.fore_color.rgb = bg
        else:
            row.fill.background()
        rtf = row.text_frame; rtf.word_wrap = False
        rtf.margin_top = rtf.margin_bottom = 0; rtf.margin_left = Pt(6); rtf.margin_right = 0
        rtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        rp = rtf.paragraphs[0]; rp.alignment = PP_ALIGN.LEFT
        _run(rp, ln if ln else " ", 10.5, fg, mono=True)
        y += lh
    if more:
        row = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx + Inches(0.06), y, cw - Inches(0.12), lh)
        no_line(row); row.fill.background()
        mp = row.text_frame.paragraphs[0]; mp.alignment = PP_ALIGN.LEFT
        _run(mp, f"… (+{more} more lines)", 10.5, SUBT, mono=True)
    footer(s, accent, name, n, total)


def add_profile(prs, e, accent, name, n, total):
    s = blank(prs); header(s, e, "Profile", accent)
    summ = e.get("profile_summary") or ""; mets = e.get("profile_metrics") or []
    bullets = (e.get("profile_bullets") or [])[:7]
    if summ or bullets:
        tf = textbox(s, Inches(0.7), Inches(1.85), Inches(7.2 if mets else 11.9), Inches(5.0),
                     MSO_ANCHOR.MIDDLE)
        # summary prose (if any) then the finding bullets — show BOTH for richer slides
        n_items = (1 if summ else 0) + len(bullets)
        sz = 19 if n_items <= 3 else (16 if n_items <= 5 else 14)
        first = True
        if summ:
            p = tf.paragraphs[0]; p.line_spacing = 1.3; p.space_after = Pt(8)
            add_rich(p, summ, sz + 1, INK, accent); first = False
        for b in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            p.line_spacing = 1.16; p.space_after = Pt(7)
            _run(p, "•  ", sz, accent, bold=True)
            add_rich(p, b, sz, INK, accent)
    if mets:
        rows = len(mets)
        tbl = s.shapes.add_table(rows, 2, Inches(8.2), Inches(1.85),
                                 Inches(4.45), Inches(min(5.0, 0.4 * rows))).table
        tbl.first_row = False; tbl.horz_banding = False
        tbl.columns[0].width = Inches(2.95); tbl.columns[1].width = Inches(1.5)
        for i, (k, v) in enumerate(mets):
            for j, val in enumerate((k, v)):
                c = tbl.cell(i, j); c.margin_top = c.margin_bottom = Pt(2)
                c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else CODEBG
                pr = c.text_frame.paragraphs[0]
                _run(pr, val, 11.5, MUTE if j == 0 else INK, bold=(j == 1))
                if j == 1:
                    pr.alignment = PP_ALIGN.RIGHT
    footer(s, accent, name, n, total)


def add_result(prs, e, gif, accent, name, n, total):
    s = blank(prs); header(s, e, "Result", accent)
    run = f"this run: {e['mfu']:.1f}% MFU" if e.get("mfu") is not None else "this run: —"
    if e.get("tps"):
        run += f" · {e['tps']:,} tok/s/chip"
    v = (e.get("verdict") or "").lower()
    tf = textbox(s, Inches(0.6), Inches(1.6), Inches(9), Inches(0.4))
    _run(tf.paragraphs[0], run, 13, SUBT)
    if v in VBADGE:
        fg, bg = VBADGE[v]
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.0), Inches(1.58), Inches(1.7), Inches(0.42))
        badge.fill.solid(); badge.fill.fore_color.rgb = bg; no_line(badge); badge.adjustments[0] = 0.3
        pb = badge.text_frame.paragraphs[0]; pb.alignment = PP_ALIGN.CENTER
        _run(pb, v, 12, fg, bold=True)
    s.shapes.add_picture(gif, Inches(1.55), Inches(2.15), height=Inches(4.75))
    footer(s, accent, name, n, total)


def add_verdict(prs, e, accent, name, n, total):
    s = blank(prs); header(s, e, "Verdict", accent)
    v = (e.get("verdict") or "").lower(); top = Inches(1.9)
    if v in VBADGE:
        fg, bg = VBADGE[v]
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), top, Inches(2.0), Inches(0.5))
        badge.fill.solid(); badge.fill.fore_color.rgb = bg; no_line(badge); badge.adjustments[0] = 0.3
        pb = badge.text_frame.paragraphs[0]; pb.alignment = PP_ALIGN.CENTER
        _run(pb, v.upper(), 13, fg, bold=True)
        top = Inches(2.7)
    tf = textbox(s, Inches(0.7), top, Inches(11.9), Inches(4.2))
    first = True
    for para in (e.get("verdict_paras") or []):
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_after = Pt(12); p.line_spacing = 1.32
        add_rich(p, para, 19, INK, accent)
    footer(s, accent, name, n, total)


def title_slide(prs, agent_name, accent, n, ctx_label, best_by_seq, seqs):
    s = blank(prs)
    tx = Inches(0.9)
    if ICON_PNG:
        pic = s.shapes.add_picture(ICON_PNG, Inches(0.9), Inches(1.28), height=Inches(0.82))
        tx = Inches(0.9) + pic.width + Inches(0.22)
    tf = textbox(s, tx, Inches(1.5), Inches(13.0) - tx, Inches(0.6))
    _run(tf.paragraphs[0], f"{agent_name} · autonomous TPU optimization", 16, accent, bold=True)
    tf = textbox(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(2.0))
    _run(tf.paragraphs[0], "Autonomous TPU Optimization", 44, INK, bold=True)
    p = tf.add_paragraph(); _run(p, f"Qwen3-8B · v6e-8 · {ctx_label}", 22, MUTE)
    tf = textbox(s, Inches(0.95), Inches(4.6), Inches(11.5), Inches(2))
    p = tf.paragraphs[0]; p.line_spacing = 1.3
    _run(p, f"{n} experiments, run autonomously — a hypothesis it formed, a change it wrote, "
            f"a profile it read, a verdict it reached, then the next idea.", 17, MUTE)
    for s_ in seqs:
        bb = best_by_seq.get(s_)
        if bb is None:
            continue
        p = tf.add_paragraph(); p.space_before = Pt(12)
        if len(seqs) > 1:
            _run(p, f"{SEQ_LABEL.get(s_, s_)}  ", 18, MUTE, bold=True)
        _run(p, "peak MFU  ", 19, INK, bold=True)
        _run(p, f"{bb:.1f}%", 19, accent, bold=True)
        tv = bd.MAXTEXT_MFU.get(s_)
        if tv:
            _run(p, f"     (MaxText SOTA {tv:.1f}%)", 15, MUTE)


def climb_slide(prs, gif, accent, agent_name):
    s = blank(prs)
    if ICON_PNG:
        s.shapes.add_picture(ICON_PNG, Inches(0.6), Inches(0.4), height=Inches(0.66))
    tf = textbox(s, Inches(1.4), Inches(0.36), Inches(11.4), Inches(0.9))
    _run(tf.paragraphs[0], "MFU frontier — the autonomous climb", 24, INK, bold=True)
    p = tf.add_paragraph()
    _run(p, f"Result for {agent_name} · Qwen3-8B on TPU v6e-8", 14, accent, bold=True)
    s.shapes.add_picture(gif, Inches(2.55), Inches(1.55), height=Inches(5.05))


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--series", default="qwen3-ag-jax")
    ap.add_argument("--seq", default="2048")
    ap.add_argument("--diffs-dir", default=os.path.join(SCRIPT_DIR, "diffs"))
    ap.add_argument("--out", default=os.path.join(SCRIPT_DIR, "deck.pptx"))
    ap.add_argument("--cap", type=int, default=70,
                    help="if metric experiments exceed this, keep only supported/frontier/baseline")
    args = ap.parse_args()

    model, agent, lane = bd.parse_series(args.series)
    man, _ = bd.load_manifest(args.series, model, agent, lane, args.diffs_dir)
    bd.enrich_from_pages(man["experiments"], model, agent, lane)
    exps = man["experiments"]
    if args.seq == "all":
        exps = [e for e in exps if e.get("seq") in (2048, 8192)]     # drop context-less runs
    else:
        exps = [e for e in exps if e.get("seq") == int(args.seq)]
    exps = [e for e in exps if e.get("mfu") is not None or e.get("tps") is not None]
    exps = bd.compute_frontier(exps)
    # large lanes (e.g. cx has 300+): keep the meaningful climb — wins + frontier-advancers + baselines
    capped = None
    if len(exps) > args.cap:
        capped = len(exps)
        sup = lambda e: (e.get("verdict") or "").lower() in ("supported", "baseline")
        keep = [e for e in exps if e.get("isNewBest") or sup(e)]        # wins + frontier + baselines
        if len(keep) > args.cap:                                        # still huge → frontier + baselines only
            keep = [e for e in exps if e.get("isNewBest")
                    or (e.get("verdict") or "").lower() == "baseline"]
        exps = bd.compute_frontier(keep)

    accent = man.get("color", AGENTS[agent]["color"]); accent_rgb = hexrgb(accent)
    name = man.get("agent_name", AGENTS[agent]["name"])
    seqs = sorted({e["seq"] for e in exps if e.get("seq") in (2048, 8192)})
    targets = {s: bd.MAXTEXT_MFU[s] for s in seqs if s in bd.MAXTEXT_MFU}
    best_by_seq = {s: _seq_frontier(exps, len(exps) - 1, s)[0] for s in seqs}
    ctx_label = " + ".join(SEQ_LABEL.get(s, str(s)) for s in seqs) + " context"
    N = len(exps)

    tmp = tempfile.mkdtemp()
    global ICON_PNG
    ICON_PNG = agent_icon_png(agent, tmp)
    # base frames (cumulative chart at each experiment) → reused for the summary climb GIF
    frame_imgs = []
    for k in range(N):
        p = os.path.join(tmp, f"fr{k:03d}.png")
        draw_chart(exps, k, p, accent, targets, seqs)
        frame_imgs.append(Image.open(p).convert("RGB").quantize(colors=48, method=Image.MEDIANCUT))

    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    title_slide(prs, name, accent_rgb, N, ctx_label, best_by_seq, seqs)
    for i, e in enumerate(exps):
        isbase = (e.get("verdict") or "").lower() == "baseline"
        if not isbase and (e.get("hypothesis_statement") or e.get("hypothesis_title")):
            add_hyp(prs, e, accent_rgb, name, i + 1, N)
        if not isbase and (e.get("diff") or "").strip():
            add_change(prs, e, accent_rgb, name, i + 1, N)
        if not isbase and (e.get("profile_summary") or e.get("profile_metrics") or e.get("profile_bullets")):
            add_profile(prs, e, accent_rgb, name, i + 1, N)
        gif = os.path.join(tmp, f"res{i:03d}.gif")
        build_transition_gif(exps, i, gif, accent, targets, seqs, tmp)   # single move: prev → this exp
        add_result(prs, e, gif, accent_rgb, name, i + 1, N)
        if not isbase and e.get("verdict_paras"):
            add_verdict(prs, e, accent_rgb, name, i + 1, N)

    full = os.path.join(tmp, "climb.gif")
    make_gif(frame_imgs, N - 1, full, hold=8, dur=470)
    climb_slide(prs, full, accent_rgb, name)

    prs.save(args.out)
    sz = os.path.getsize(args.out)
    peak = " / ".join(f"{SEQ_LABEL.get(s, s)} {best_by_seq[s]:.1f}%" for s in seqs if best_by_seq.get(s))
    capnote = f" (capped from {capped})" if capped else ""
    print(f"Wrote {args.out}  ({len(prs.slides)} slides, {N} experiments{capnote}, "
          f"contexts={ctx_label}, peak {peak}, {sz//1024}KB)")


if __name__ == "__main__":
    main()
