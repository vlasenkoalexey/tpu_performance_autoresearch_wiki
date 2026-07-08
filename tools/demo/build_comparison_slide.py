#!/usr/bin/env python3
"""One-slide cross-agent comparison: all 4 agents' best-MFU hill-climb, split 2k vs 8k.

    python tools/demo/build_comparison_slide.py --out /path/compare.pptx

Two panels (2k, 8k); each plots every agent's running-best MFU frontier over optimization
progress, in brand-aligned colours (Claude brown/orange, Antigravity blue, Codex teal,
Fable5 red). Legend uses the agent logos (Claude rendered in its brown/orange brand colour,
not the default black sunburst). Data from wiki/analyses/qwen3/mfu_data.json.
"""
import argparse, json, os, tempfile
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
DATA = os.path.join(SCRIPT_DIR, "..", "..", "wiki", "analyses", "qwen3", "mfu_data.json")

AGENTS = ["cc", "ag", "cx", "cc5"]
NAME = {"cc": "Claude Opus 4.8", "ag": "Antigravity Gemini 3.1 Pro",
        "cx": "Codex GPT-5.5", "cc5": "Claude Fable5"}
ICON_SRC = {"cc": "cc", "cc5": "cc", "cx": "cx", "ag": "ag"}   # both Claude models use the Claude mark
# brand-aligned colours (Claude = brown/orange per request; Antigravity = its orange; Codex = OpenAI teal)
COLOR = {"cc": "#A8552E", "ag": "#ff7f0e", "cx": "#10A37F", "cc5": "#D62728"}
MAXTEXT = {2048: 36.6, 8192: 39.8}
INK, MUTE, SUBT, LINE = "#1a1d21", "#5b6570", "#8a929c", "#e6e9ee"
SW, SH = Inches(13.333), Inches(7.5)


def hexrgb(h):
    h = h.lstrip("#"); return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def series(rows):
    """Per-experiment points in run order: x=experiment #, mfu, fail (miss), best (running best)."""
    rows = sorted([r for r in rows if r.get("mfu") is not None], key=lambda r: r["order"])
    pts, b = [], None
    for i, r in enumerate(rows):
        v = (r.get("verdict") or "").lower()
        if v != "invalid":
            b = r["mfu"] if (b is None or r["mfu"] > b) else b
        pts.append({"x": i + 1, "mfu": r["mfu"],
                    "fail": v in ("invalid", "refuted", "inconclusive"), "best": b})
    return pts


def render_icon(agent, tmp):
    """Legend logo tinted to the agent's colour. Both Claude models use the Claude mark;
    Codex the OpenAI mark; Antigravity its own (natural). No fabricated marks."""
    import resvg_py
    src = ICON_SRC[agent]
    svg = os.path.join(SCRIPT_DIR, "assets", f"{src}.svg")
    out = os.path.join(tmp, f"leg_{agent}.png")
    s = open(svg, encoding="utf-8").read()
    if src in ("cc", "cx"):                      # monochrome (currentColor) → tint to line colour
        s = s.replace("currentColor", COLOR[agent])
    open(out, "wb").write(bytes(resvg_py.svg_to_bytes(svg_string=s, width=256)))
    return out


def draw_frame(p, data, maxN, ylim, path, fig_wh):
    """One animation frame at reveal fraction p: each agent's points+misses+frontier
    revealed up to p·(its experiment count). Style matches the per-experiment deck charts
    (filled = hit, hollow = miss) with the running-best frontier line."""
    fig, axes = plt.subplots(1, 2, figsize=fig_wh, dpi=140)
    fig.patch.set_facecolor("white")
    for ax, seq in zip(axes, (2048, 8192)):
        ax.set_facecolor("white")
        tv = MAXTEXT[seq]
        ax.axhline(tv, ls=(0, (6, 5)), color="#9aa0a6", lw=1.4, zorder=2)
        ax.text(0.985, tv, f"MaxText SOTA {tv:.1f}%", ha="right", va="bottom", color="#8a929c",
                fontsize=10, transform=ax.get_yaxis_transform())
        for ag in AGENTS:
            pts = data[(ag, seq)]
            if not pts:
                continue
            k = max(1, round(p * len(pts))); rev = pts[:k]
            fx = [q["x"] for q in rev if q["best"] is not None]
            fy = [q["best"] for q in rev if q["best"] is not None]
            if len(fx) > 1:
                ax.step(fx, fy, where="post", color=COLOR[ag], lw=2.6, zorder=5, solid_capstyle="round")
            miss = [(q["x"], q["mfu"]) for q in rev if q["fail"]]
            hit = [(q["x"], q["mfu"]) for q in rev if not q["fail"]]
            if miss:
                ax.scatter(*zip(*miss), facecolors="white", edgecolors=COLOR[ag], s=15,
                           linewidths=0.8, zorder=4)
            if hit:
                ax.scatter(*zip(*hit), color=COLOR[ag], s=16, linewidths=0, zorder=6)
            if fx:
                ax.annotate(f"{fy[-1]:.1f}%", (fx[-1], fy[-1]), textcoords="offset points",
                            xytext=(5, 2), fontsize=10, fontweight="bold", color=COLOR[ag])
        ax.set_title(f"{seq // 1024}k context", fontsize=15, fontweight="bold", color=INK, pad=8)
        ax.set_xlim(0.5, maxN[seq] * 1.06); ax.set_ylim(*ylim[seq])
        ax.set_xlabel("experiment #", fontsize=11, color=MUTE)
        ax.set_ylabel("best MFU (%, causal)", fontsize=11, color=MUTE)
        ax.grid(True, axis="y", color="#eceff3", lw=1)
        for sp in ("top", "right"):
            ax.spines[sp].set_visible(False)
        for sp in ("left", "bottom"):
            ax.spines[sp].set_color("#d7dbe0")
        ax.tick_params(colors="#8a929c", labelsize=9)
    fig.tight_layout(pad=1.4)
    fig.savefig(path, facecolor="white"); plt.close(fig)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=os.path.join(SCRIPT_DIR, "compare.pptx"))
    ap.add_argument("--frames", type=int, default=32)
    args = ap.parse_args()
    rows = json.load(open(DATA))
    data, best = {}, {}
    for ag in AGENTS:
        for seq in (2048, 8192):
            pts = series([r for r in rows if r["lane"] == ag and r["seq"] == seq])
            data[(ag, seq)] = pts
            best[(ag, seq)] = pts[-1]["best"] if pts else None
    maxN = {seq: max((len(data[(ag, seq)]) for ag in AGENTS), default=1) for seq in (2048, 8192)}
    ylim = {}
    for seq in (2048, 8192):
        allm = [q["mfu"] for ag in AGENTS for q in data[(ag, seq)] if q["mfu"] is not None] + [MAXTEXT[seq]]
        ylim[seq] = (max(0, min(allm) * 0.85), max(allm) * 1.12)

    tmp = tempfile.mkdtemp()
    fig_wh = (12.6, 5.15)                       # slightly larger diagram
    F = args.frames; frames = []
    for i in range(F):
        fp = os.path.join(tmp, f"f{i:03d}.png")
        draw_frame((i + 1) / F, data, maxN, ylim, fp, fig_wh)
        frames.append(Image.open(fp).convert("RGB").quantize(colors=64, method=Image.MEDIANCUT))
    gif = os.path.join(tmp, "climb.gif")
    durs = [120] * F; durs[0] = 550; durs[-1] = 2000
    frames[0].save(gif, save_all=True, append_images=frames[1:], duration=durs, loop=0, optimize=True)
    icons = {ag: render_icon(ag, tmp) for ag in AGENTS}

    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = hexrgb("#ffffff")
    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.95)).text_frame
    tb.word_wrap = True
    r = tb.paragraphs[0].add_run(); r.text = "Four agents, one problem — the autonomous MFU climb"
    r.font.size = Pt(25); r.font.bold = True; r.font.name = "Calibri"; r.font.color.rgb = hexrgb(INK)
    p = tb.add_paragraph(); r = p.add_run()
    r.text = ("Qwen3-8B · TPU v6e-8 · best MFU (causal) vs experiment # — filled = hit, "
              "hollow = miss, line = running best, dashed = MaxText SOTA")
    r.font.size = Pt(12.5); r.font.name = "Calibri"; r.font.color.rgb = hexrgb(MUTE)
    s.shapes.add_picture(gif, Inches(0.4), Inches(1.28), width=Inches(12.55))

    # icon legend row (full model names)
    x0, y, step = Inches(0.55), Inches(6.62), Inches(3.02)
    for i, ag in enumerate(AGENTS):
        x = x0 + step * i
        s.shapes.add_picture(icons[ag], x, y, height=Inches(0.44))
        tf = s.shapes.add_textbox(x + Inches(0.54), y - Inches(0.06), Inches(2.42), Inches(0.7)).text_frame
        tf.word_wrap = True
        rn = tf.paragraphs[0].add_run(); rn.text = NAME[ag]
        rn.font.size = Pt(12.5); rn.font.bold = True; rn.font.name = "Calibri"; rn.font.color.rgb = hexrgb(COLOR[ag])
        pp = tf.add_paragraph(); rn = pp.add_run()
        b2, b8 = best[(ag, 2048)], best[(ag, 8192)]
        rn.text = f"2k {b2:.1f}%" + (f" · 8k {b8:.1f}%" if b8 else "")
        rn.font.size = Pt(10.5); rn.font.name = "Calibri"; rn.font.color.rgb = hexrgb(MUTE)

    prs.save(args.out)
    print(f"Wrote {args.out} ({os.path.getsize(args.out)//1024}KB)")
    for ag in AGENTS:
        print(f"  {NAME[ag]:26} 2k={best[(ag,2048)]}  8k={best[(ag,8192)]}  "
              f"(n2k={len(data[(ag,2048)])}, n8k={len(data[(ag,8192)])})")


if __name__ == "__main__":
    main()
